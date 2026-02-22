# danny_toolkit/core/doc_loader.py — Document loader for RAG indexing
"""
Ondersteunde formaten:
  - PDF  (.pdf)   — pdfplumber met tabel-extractie + fallback PyPDF2
  - Word (.docx)  — python-docx
  - PowerPoint (.pptx) — python-pptx
  - Excel (.xlsx) — openpyxl
  - EPUB (.epub)  — ebooklib + BeautifulSoup
  - Tekst (.txt, .md, .py, .json, .csv, .log, .html, .xml, .rst)

Chunks bevatten metadata: bron, pagina/sectie, chunk nummer.
"""

import logging
import re
from pathlib import Path

logger = logging.getLogger(__name__)


def load_directory(directory: str, chunk_size: int = 500, overlap: int = 50) -> list[dict]:
    """Load all supported files from a directory (or single file) and split into chunks.

    Returns list of dicts: {"text": str, "source": str, "chunk": int, "page": int|None}
    """
    target = Path(directory)
    if not target.exists():
        raise FileNotFoundError(f"Pad niet gevonden: {target}")

    loaders = {
        ".pdf": _load_pdf,
        ".docx": _load_docx,
        ".pptx": _load_pptx,
        ".xlsx": _load_xlsx,
        ".epub": _load_epub,
        ".txt": _load_text,
        ".md": _load_text,
        ".py": _load_text,
        ".json": _load_text,
        ".csv": _load_text,
        ".log": _load_text,
        ".html": _load_html,
        ".xml": _load_text,
        ".rst": _load_text,
    }

    chunks = []
    file_count = 0
    skipped = []

    # Enkel bestand of directory
    if target.is_file():
        files = [target]
    else:
        files = sorted(target.rglob("*"))

    for filepath in files:
        if not filepath.is_file():
            continue
        ext = filepath.suffix.lower()
        loader = loaders.get(ext)
        if loader is None:
            continue

        try:
            result = loader(filepath)
        except Exception as e:
            print(f"  [!] Fout bij {filepath.name}: {e}")
            skipped.append(filepath.name)
            continue

        # result is list of {"text": str, "page": int|None}
        if not result:
            continue

        file_count += 1
        for page_block in result:
            text = page_block["text"]
            page = page_block.get("page")
            if not text or not text.strip():
                continue

            file_chunks = _chunk_text(text, chunk_size, overlap)
            for i, chunk in enumerate(file_chunks):
                chunks.append({
                    "text": chunk,
                    "source": str(filepath),
                    "chunk": i,
                    "page": page,
                })

    print(f"  {file_count} bestanden geladen, {len(chunks)} chunks gemaakt")
    if skipped:
        print(f"  [!] {len(skipped)} bestanden overgeslagen: {', '.join(skipped)}")
    return chunks


# ═══════════════════════════════════════════
# PDF — pdfplumber (met tabel-extractie)
# ═══════════════════════════════════════════

def _load_pdf(filepath: Path) -> list[dict]:
    """Load PDF with pdfplumber (tables + text), fallback to PyPDF2."""
    try:
        return _load_pdf_plumber(filepath)
    except Exception as e:
        logger.debug("pdfplumber mislukt, fallback naar PyPDF2: %s", e)
        return _load_pdf_pypdf2(filepath)


def _load_pdf_plumber(filepath: Path) -> list[dict]:
    import pdfplumber

    pages = []
    with pdfplumber.open(str(filepath)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            parts = []

            # Tekst extractie
            text = page.extract_text()
            if text:
                parts.append(_clean_pdf_text(text))

            # Tabel extractie
            tables = page.extract_tables()
            for table in tables:
                table_text = _table_to_text(table)
                if table_text:
                    parts.append(table_text)

            combined = "\n\n".join(parts)
            if combined.strip():
                pages.append({"text": combined, "page": i})

    if not pages:
        raise ValueError("Geen tekst gevonden met pdfplumber")
    return pages


def _load_pdf_pypdf2(filepath: Path) -> list[dict]:
    from PyPDF2 import PdfReader

    reader = PdfReader(str(filepath))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            pages.append({"text": _clean_pdf_text(text), "page": i})
    return pages


def _clean_pdf_text(text: str) -> str:
    """Clean up common PDF extraction artifacts."""
    # Fix broken words (hy- phenation)
    text = re.sub(r'(\w)-\s*\n\s*(\w)', r'\1\2', text)
    # Normalize whitespace
    text = re.sub(r'[ \t]+', ' ', text)
    # Remove excessive newlines
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def _table_to_text(table: list) -> str:
    """Convert pdfplumber table to readable text."""
    if not table:
        return ""
    rows = []
    for row in table:
        cells = [str(c).strip() if c else "" for c in row]
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)


# ═══════════════════════════════════════════
# DOCX — Word documenten
# ═══════════════════════════════════════════

def _load_docx(filepath: Path) -> list[dict]:
    from docx import Document

    doc = Document(str(filepath))
    parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # Tabellen uit Word
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append("\n".join(rows))

    combined = "\n\n".join(parts)
    if combined.strip():
        return [{"text": combined, "page": None}]
    return []


# ═══════════════════════════════════════════
# XLSX — Excel spreadsheets
# ═══════════════════════════════════════════

def _load_xlsx(filepath: Path) -> list[dict]:
    from openpyxl import load_workbook

    wb = load_workbook(str(filepath), read_only=True, data_only=True)
    sheets = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            text = f"Sheet: {sheet_name}\n" + "\n".join(rows)
            sheets.append({"text": text, "page": None})

    wb.close()
    return sheets


# ═══════════════════════════════════════════
# PPTX — PowerPoint presentaties
# ═══════════════════════════════════════════

def _load_pptx(filepath: Path) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(str(filepath))
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    parts.append("\n".join(rows))
        combined = "\n\n".join(parts)
        if combined.strip():
            slides.append({"text": combined, "page": i})

    return slides


# ═══════════════════════════════════════════
# EPUB — E-books
# ═══════════════════════════════════════════

def _load_epub(filepath: Path) -> list[dict]:
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup

    book = epub.read_epub(str(filepath), options={"ignore_ncx": True})
    chapters = []

    for i, item in enumerate(book.get_items_of_type(ebooklib.ITEM_DOCUMENT), 1):
        soup = BeautifulSoup(item.get_content(), "html.parser")

        for tag in soup(["script", "style"]):
            tag.decompose()

        text = soup.get_text(separator="\n", strip=True)
        text = re.sub(r"\n{3,}", "\n\n", text)
        text = re.sub(r"[ \t]+", " ", text)

        if text.strip():
            chapters.append({"text": text.strip(), "page": i})

    return chapters


# ═══════════════════════════════════════════
# Tekst bestanden
# ═══════════════════════════════════════════

def _load_text(filepath: Path) -> list[dict]:
    text = filepath.read_text(encoding="utf-8", errors="ignore")
    if text.strip():
        return [{"text": text, "page": None}]
    return []


# ═══════════════════════════════════════════
# HTML bestanden
# ═══════════════════════════════════════════

def _load_html(filepath: Path) -> list[dict]:
    from bs4 import BeautifulSoup

    html = filepath.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")

    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()

    text = soup.get_text(separator="\n", strip=True)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)

    if text.strip():
        return [{"text": text.strip(), "page": None}]
    return []


# ═══════════════════════════════════════════
# Chunking
# ═══════════════════════════════════════════

def _chunk_text(text: str, chunk_size: int, overlap: int) -> list[str]:
    """Split text into overlapping chunks by word count."""
    words = text.split()
    if len(words) <= chunk_size:
        return [text.strip()] if text.strip() else []

    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        if chunk.strip():
            chunks.append(chunk.strip())
        start += chunk_size - overlap

    return chunks
